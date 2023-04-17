import streamlit as st
import openai as ai
import numpy as np
import pandas as pd
# import aspose.words as aw
from htmldocx import HtmlToDocx
from bing_image_downloader import downloader
import os
from PIL import Image
import base64
import os
import json
import pickle
import uuid
import re
import shutil
#from streamlit_login_auth_ui.widgets import __login__
import requests
from google.oauth2 import service_account
# from gsheetsdb import connect
from gspread_pandas import Spread,Client
import gspread_pandas
from datetime import datetime
from pptx import Presentation
from io import BytesIO
from pptx.util import Inches, Pt
import streamlit_analytics

# Config

ai.api_key = st.secrets["openai_api_key"]

API_KEY = st.secrets["openai_api_key"]
API_ENDPOINT = "https://api.openai.com/v1/chat/completions"

# For GPT-4
def generate_chat_completion(messages, model="gpt-4", temperature=1, max_tokens=None):
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {API_KEY}",
    }

    data = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
    }

    if max_tokens is not None:
        data["max_tokens"] = max_tokens

    response = requests.post(API_ENDPOINT, headers=headers, data=json.dumps(data))

    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        raise Exception(f"Error {response.status_code}: {response.text}")

# For Da-Vinci 3
def generate_response(MODEL, PROMPT, MAX_TOKENS=750, TEMP=0.99, TOP_P=0.5, N=1, FREQ_PEN=0.3, PRES_PEN = 0.9):
  response = ai.Completion.create(
          engine = MODEL,
          # engine="text-davinci-002", # OpenAI has made four text completion engines available, named davinci, ada, babbage and curie. We are using davinci, which is the most capable of the four.
          prompt=PROMPT, # The text file we use as input (step 3)
          max_tokens=MAX_TOKENS, # how many maximum characters the text will consists of.
          temperature=TEMP,
          # temperature=int(temperature), # a number between 0 and 1 that determines how many creative risks the engine takes when generating text.,
          top_p=TOP_P, # an alternative way to control the originality and creativity of the generated text.
          n=N, # number of predictions to generate
          frequency_penalty=FREQ_PEN, # a number between 0 and 1. The higher this value the model will make a bigger effort in not repeating itself.
          presence_penalty=PRES_PEN # a number between 0 and 1. The higher this value the model will make a bigger effort in talking about new topics.
      )
  return response['choices'][0]['text']

MODEL = 'text-davinci-003'
streamlit_analytics.start_tracking()
st.header('Knowledge Organiser Generator')
st.sidebar.image('pedagogical_18.png')
st.sidebar.markdown("This worksheet generator was created using OpenAI's generative AI. Please use it carefully and check any output before using it with learners as it could be biased or wrong. ")
st.markdown("Other Pedagogical apps to check out: [worksheet generator](https://pedagogical.app/)")

txt_button = st.checkbox('What is a knowledge organiser?', help='The learning science behind Knowledge Organisers')
if txt_button:
    st.subheader('Key Features')
    overview_txt = st.markdown("""
    - [Precisely specifying what students should know by the end of the unit.](https://improvingteaching.co.uk/2017/04/23/better-planning-better-teaching-better-learning-a-template/)
    - Contain knowledge which is [broken down into manageble units.](https://www.supermemo.com/en/blog/twenty-rules-of-formulating-knowledge)
    - Knowledge organisers should be structured so students can use them to self-quiz.(https://researchschool.org.uk/bradford/news/knowledge-organisers-facilitating-retrieval)
    - Students should be taught how to use them to self-quiz. """)

    st.subheader('Tips for Use')
    key_concepts_txt = st.markdown('''
    - Use for [self-checking after a task.](https://classteaching.wordpress.com/2018/09/14/using-knowledge-organisers-to-improve-retrieval-practice/)
    - [Don't test all at once.](https://classteaching.wordpress.com/2018/09/14/using-knowledge-organisers-to-improve-retrieval-practice/)
    - [Students should understand concepts, learn how the concepts relate to others and form an overall picture and, finally, use a knoweldge organiser to memorise key information.](https://www.supermemo.com/en/blog/twenty-rules-of-formulating-knowledge)
    ''')

    st.subheader('Critiques of Knowledge Organisers')
    lethal_injections_txt = st.markdown('''
    - Knowledge Organisers were originally developed at Michaela School, and pioneered by Joe Kirby, according to [this podcast](https://www.ollielovell.com/errr/harryfletcherwood2/) Michaela School has reduced their use due to their restrictive nature..
    - Some teachers find they lead to students ['blindly copying answers'](https://shallteach.wordpress.com/2020/07/25/knowledge-organisers-a-failed-revolution/).
    - Knowledge organisers may not work as well for [subjects with more procedural knowledge](https://tothereal.wordpress.com/2018/06/04/why-maths-teachers-dont-like-knowledge-organisers/), e.g. Maths.
    - They present knoweldge as linear, when [information is often relational.](https://shallteach.wordpress.com/2020/07/25/knowledge-organisers-a-failed-revolution/)
    ''')

email = st.text_input('Email')
topic = st.text_input('Knowledge Organiser Topic')
reading_age = st.slider('Reading Age', 0, 18)

components = []
cloze_components = []
blank_components = []

st.markdown("Humanities and Sciences Options")
key_words_check = st.checkbox('Key Words')
key_concepts_check = st.checkbox('Key concepts')
timeline_check = st.checkbox('Timeline')

st.markdown("English Options")

characters_check = st.checkbox('Key Characters')
characters_quotes_check = st.checkbox('Character Quotes')
dramatic_devices_check = st.checkbox('Dramatic Devices')
plot_check = st.checkbox('Plot')

self_quiz_format = st.radio('Self-Quizzing format', ['Blank', 'No Self Quizzing'])

generate_worksheet = st.button('Generate Knowledge Organiser')

if generate_worksheet:
    with st.spinner(text="Your worksheet is in the oven 🧠 ... If you want to work with Pedagogical to improve the app please click [here](https://forms.gle/jDy1WNgrnCTWsDG16) ... Thank you!"):

        key_words_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of key words and their definitions so that student can understand the topic."
        if key_words_check:
            key_words_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": key_words_prompt}
            ]
            key_words = generate_response(MODEL, key_words_prompt) #generate_chat_completion(key_words_messages)
            components.append(key_words)
            if self_quiz_format == 'Blank':
                key_words_blank_prompt = f"Pick out the key words only from these key words and definitions: {key_words}. Leave two spaces below each."
                key_words_blank = generate_response(MODEL, key_words_blank_prompt)
                blank_components.append(key_words_blank)
            elif self_quiz_format == 'Cloze':
                key_words_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn these key words and definitions into a cloze exercise: {key_words}. Output all missing words at the end."
                key_words_cloze = generate_response(MODEL, key_words_cloze_prompt)
                cloze_components.append(key_words_cloze)


        key_concepts_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of key concepts and their definitions that would be important for them to understand this topic."
        if key_concepts_check:
            key_concepts_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": key_concepts_prompt}
            ]
            key_concepts = generate_response(MODEL, key_concepts_prompt)
            components.append(key_concepts)
            if self_quiz_format == 'Blank':
                key_concepts_blank_prompt = f"Pick out the key concepts only from these key words and definitions: {key_concepts}. Leave two spaces below each."
                key_concepts_blank = generate_response(MODEL, key_concepts_blank_prompt)
                blank_components.append(key_concepts_blank)
            elif self_quiz_format == 'Cloze':
                key_concepts_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn these key concepts and definitions into a cloze exercise: {key_concepts} Output all missing words at the end."
                key_concepts_cloze = generate_response(MODEL, key_concepts_cloze_prompt)
                cloze_components.append(key_concepts_cloze)



        timeline_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a timeline so that your student can understand the topic."
        if timeline_check:
            timeline_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": timeline_prompt}
            ]
            timeline = generate_response(MODEL, timeline_prompt)
            components.append(timeline)
            if self_quiz_format == 'Blank':
                timeline_blank_prompt = f"Pick out the dates only from this timeline: {timeline}. Leave two spaces below each."
                timeline_blank = generate_response(MODEL, timeline_blank_prompt)
                blank_components.append(timeline_blank)
            elif self_quiz_format == 'Cloze':
                timeline_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn this timeline into a cloze exercise: {timeline}. Output all missing words at the end."
                timeline_cloze = generate_response(MODEL, timeline_cloze_prompt)
                cloze_components.append(timeline_cloze)


        characters_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of important characters and their descriptions so that student can understand the topic."
        if characters_check:
            characters_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": characters_prompt}
            ]
            characters = generate_response(MODEL, characters_prompt)
            components.append(characters)
            if self_quiz_format == 'Blank':
                characters_blank_prompt = f"Pick out the characters only from this list of characters and descriptions {characters}. Leave Two spaces below each."
                characters_blank = generate_response(MODEL, characters_blank_prompt)
                blank_components.append(characters_blank)
            elif self_quiz_format == 'Cloze':
                characters_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn this list of characters and descriptions into a cloze exercise: {characters}. Output all missing words at the end."
                characters_cloze = generate_response(MODEL, characters_cloze_prompt)
                cloze_components.append(characters_cloze)






        characters_quotes_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of important quotes by characters from {topic} so that your student can understand the topic."
        if characters_quotes_check:
            characters_quotes_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": characters_quotes_prompt}
            ]
            characters_quotes = generate_response(MODEL, characters_quotes_prompt)
            components.append(characters_quotes)
            if self_quiz_format == 'Blank':
                characters_quotes_blank_prompt = f"Pick out the first two words of each of the  characters quotes only from this list of characters quotes and leave enough space after the two words so someone could finish the quote: {characters_quotes}"
                characters_quotes_blank = generate_response(MODEL, characters_quotes_blank_prompt)
                blank_components.append(characters_quotes_blank)
            elif self_quiz_format == 'Cloze':
                characters_quotes_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn this list of quotes into a cloze exercise: {characters_quotes}. Leave the missing words at the end."
                characters_quotes_cloze = generate_response(MODEL, characters_quotes_cloze_prompt)
                cloze_components.append(characters_quotes_cloze)


        dramatic_devices_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of important dramatic devices from {topic} and how they used in {topic}."
        if dramatic_devices_check:
            dramatic_devices_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": dramatic_devices_prompt}
            ]
            dramatic_devices = generate_response(MODEL, dramatic_devices_prompt )
            components.append(dramatic_devices)
            if self_quiz_format == 'Blank':
                dramatic_devices_blank_prompt = f"Pick out each numbered dramatic device only from this list of dramatic_devices {dramatic_devices}"
                dramatic_devices_blank = generate_response(MODEL, dramatic_devices_blank_prompt)
                blank_components.append(dramatic_devices_blank)
            elif self_quiz_format == 'Cloze':
                dramatic_devices_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn this list of dramatic devices into a cloze exercise: {dramatic_devices}. Leave the missing words at the end."
                dramatic_devices_cloze = generate_response(MODEL, dramatic_devices_cloze_prompt)
                cloze_components.append(dramatic_devices_cloze)

        plot_prompt = f"You are an expert teacher trying to teach your {reading_age} year old student about {topic}. Create a list of key episodes in the plot so that your student can understand the topic."
        if plot_check:
            plot_messages = [
                {"role": "system", "content": "You are an expert teacher."},
                {"role": "user", "content": plot_prompt}
            ]
            plot = generate_response(MODEL, plot_prompt)
            components.append(plot)
            if self_quiz_format == 'Blank':
                plot_blank_prompt = f"Pick out the first two words from each numbered line from the following, and leave two spaces between each pair of words: {plot}"
                plot_blank = generate_response(MODEL, plot_blank_prompt)
                blank_components.append(plot_blank)
            elif self_quiz_format == 'Cloze':
                plot_cloze_prompt = f"You are an expert teacher and you want to make a cloze exercise that removes certain words from text and print them at the end so that your students can guess the answers. Turn this list of plot parts into a cloze exercise: {plot}. Leave the missing words at the end."
                plot_cloze = generate_response(MODEL, plot_cloze_prompt)
                cloze_components.append(plot_cloze)

        components_dict = {val: name for name, val in locals().items() if val in components}
        component_names = [components_dict[val] for val in components]

        blank_components_dict = {val: name for name, val in locals().items() if val in blank_components}
        blank_component_names = [blank_components_dict[val] for val in blank_components]

        cloze_components_dict = {val: name for name, val in locals().items() if val in cloze_components}
        cloze_component_names = [cloze_components_dict[val] for val in cloze_components]
            # worksheet_end = """</body>
            # </html>"""

            # # Construct Worksheet from components
            # full_worksheet = worksheet_head

            # for component in components:
            #     if component:
            #         full_worksheet += component

            # full_worksheet += worksheet_end
    # f = open('worksheet.html','w')
    # f.write(knowledge_organiser_html)
    # f.close()

    # new_parser = HtmlToDocx()
    # new_parser.table_style = 'TableGrid'
    # new_parser.parse_html_file("worksheet.html", "worksheet")

    # file_path = 'worksheet.docx'
    # with open(file_path,"rb") as f:
    #     base64_word = base64.b64encode(f.read()).decode('utf-8')

    # with open("worksheet.docx", "rb") as word_file:
    #     wordbyte = word_file.read()


    # downloaded = st.download_button(label="Download Word Document", 
    # data=wordbyte,
    # file_name="pedagogical_worksheet.docx",
    # mime='application/octet-stream')
        

        prs = Presentation()
        title_only_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = f'{topic} Knowledge Organiser'

        cols = len(components)
        rows = 2
        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        for i, component in enumerate(components):
            # set column widths
            table.columns[i].width = Inches(2.0)
            # table.columns[1].width = Inches(4.0)

            # table.font.size = Pt(6)

            # # write column headings
            # table.cell(0, 0).text = 'Key Words'
            # table.cell(0, 1).text = 'Key Concepts'

            # # write body cells
            # table.cell(1, 0).text = key_words
            # table.cell(1, 1).text = key_concepts

            # # write column headings
            table.cell(0, i).text = component_names[i]

            # # write body cells
            table.cell(1, i).text = component

        def iter_cells(table):
            for row in table.rows:
                for cell in row.cells:
                    yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(6)

        if self_quiz_format == 'Blank':
            blank_slide_layout = prs.slide_layouts[5]
            blank_slide = prs.slides.add_slide(blank_slide_layout)
            shapes = blank_slide.shapes

            shapes.title.text = 'Quiz Yourself (Blanks)'
            cols = len(components)
            rows = 2
            left = top = Inches(2.0)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table

            for i, component in enumerate(blank_components):
                # set column widths
                table.columns[i].width = Inches(2.0)
                table.cell(0, i).text = component_names[i]
                # # write body cells
                table.cell(1, i).text = component

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(6)

       
        elif self_quiz_format == 'Cloze':
            cloze_slide_layout = prs.slide_layouts[5]
            cloze_slide = prs.slides.add_slide(cloze_slide_layout)
            shapes = cloze_slide.shapes

            shapes.title.text = 'Quiz Yourself (Cloze)'
            cols = len(components)
            rows = 2
            left = top = Inches(2.0)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table

            for i, component in enumerate(cloze_components):
                # set column widths
                table.columns[i].width = Inches(2.0)
                table.cell(0, i).text = component_names[i]
                # # write body cells
                table.cell(1, i).text = component

            for cell in iter_cells(table):
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(6)

        elif self_quiz_format == 'No Self Quizzing':
            pass

        prs.save('knowledge_organiser.pptx')

        # save the output into binary form
        binary_output = BytesIO()
        prs.save(binary_output) 

        st.download_button(label = 'Download Powerpoint',
                        data = binary_output.getvalue(),
                        file_name = 'pedagogical_knowledge_organiser.pptx')

        scope = ['https://spreadsheets.google.com/feeds']

        credentials = service_account.Credentials.from_service_account_info(
                        st.secrets["gcp_service_account"], scopes = scope)
        client = Client(scope=scope,creds=credentials)
        spreadsheetname = st.secrets["private_gsheets_knowledge_organiser_url"]
        spread = Spread(spreadsheetname,client = client)
        read_df = spread.sheet_to_df(index=False)
        emails = list(read_df.emails.values)
        prompts = list(read_df.prompts.values)
        dates = list(read_df.dates.values)

        today = datetime.now()
        emails.append(email)
        prompts.append(topic)
        dates.append(today)
        def update_the_spreadsheet(spreadsheetname,dataframe):
            spread.df_to_sheet(dataframe,sheet = spreadsheetname,index = False)
        d = {'emails': emails, 'prompts': prompts, 'dates': dates}
        df = pd.DataFrame(data=d)
        update_the_spreadsheet('Sheet1',df)

streamlit_analytics.stop_tracking()